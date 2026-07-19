"""
document_service.py — File upload, text extraction, and document management.

Supported formats:
    PDF   (.pdf)  — extracted via PyMuPDF (fitz)
    Word  (.docx) — extracted via python-docx
    Word  (.doc)  — fallback: latin-1 bytes → strip binary noise
    PPT   (.pptx) — extracted via python-pptx
    PPT   (.ppt)  — fallback text extraction (same as .doc)
    Text  (.txt)  — direct UTF-8 read

Storage layout:
    backend/uploads/<user_id>/<uuid>_<filename>

Database metadata is persisted via DBService (documents table).

Improvements over v1:
  • Text preprocessing: normalise whitespace, strip control characters,
    collapse runs of blank lines — produces cleaner LLM prompts.
  • Duplicate detection: if the same user uploads a file with the same
    filename and size, the existing document_id is returned.
  • Chunked context helper: for very large documents, returns a
    keyword-relevance-scored slice rather than a simple head truncation.
  • Cleanup helper: removes files older than N hours from disk.
"""

import os
import re
import uuid
import logging
import asyncio
import datetime
import unicodedata
from pathlib import Path
from typing import Optional

import aiofiles

from services.db_service import DBService

logger = logging.getLogger("StudyMate_API.document_service")

# ─────────────────────────────────────────────────────────────────────────────
# Constants
# ─────────────────────────────────────────────────────────────────────────────
BACKEND_DIR = Path(__file__).parent.parent
UPLOAD_DIR  = BACKEND_DIR / "uploads"

ALLOWED_EXTENSIONS  = {".pdf", ".doc", ".docx", ".ppt", ".pptx", ".txt"}
MAX_FILE_SIZE_BYTES = 20 * 1024 * 1024   # 20 MB
MAX_CONTEXT_CHARS   = 32_000             # ~10 500 tokens — safe for gpt-4o 128k window


# ─────────────────────────────────────────────────────────────────────────────
# Text extractors
# ─────────────────────────────────────────────────────────────────────────────

def _extract_pdf(file_path: Path) -> str:
    """Extract text from a PDF using PyMuPDF (fitz)."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(str(file_path))
        pages = [page.get_text() for page in doc]
        doc.close()
        return "\n".join(pages).strip()
    except ImportError:
        logger.error("PyMuPDF not installed — cannot extract PDF text.")
        return ""
    except Exception as exc:
        logger.error("PDF extraction failed for %s: %s", file_path, exc)
        return ""


def _extract_docx(file_path: Path) -> str:
    """Extract text from a .docx file."""
    try:
        from docx import Document
        doc = Document(str(file_path))
        parts: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        parts.append(cell.text.strip())
        return "\n".join(parts).strip()
    except ImportError:
        logger.error("python-docx not installed — cannot extract DOCX text.")
        return ""
    except Exception as exc:
        logger.error("DOCX extraction failed for %s: %s", file_path, exc)
        return ""


def _extract_doc(file_path: Path) -> str:
    """Best-effort extraction from a legacy .doc / .ppt binary file."""
    try:
        with open(file_path, "rb") as f:
            raw = f.read()
        text = raw.decode("latin-1", errors="ignore")
        printable = "".join(c for c in text if c.isprintable() or c in "\n\r\t")
        return printable.strip()
    except Exception as exc:
        logger.error("DOC/PPT extraction failed for %s: %s", file_path, exc)
        return ""


def _extract_pptx(file_path: Path) -> str:
    """Extract text from all slides in a .pptx file."""
    try:
        from pptx import Presentation
        prs = Presentation(str(file_path))
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_parts = [f"[Slide {i}]"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())
            slides_text.append("\n".join(slide_parts))
        return "\n\n".join(slides_text).strip()
    except ImportError:
        logger.error("python-pptx not installed — cannot extract PPTX text.")
        return ""
    except Exception as exc:
        logger.error("PPTX extraction failed for %s: %s", file_path, exc)
        return ""


def _extract_txt(file_path: Path) -> str:
    """Read a plain text file."""
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read().strip()
    except Exception as exc:
        logger.error("TXT extraction failed for %s: %s", file_path, exc)
        return ""


def extract_text(file_path: Path, file_type: str) -> str:
    """Dispatch to the correct extractor based on file extension."""
    ext = file_type.lower()
    dispatch = {
        ".pdf":  _extract_pdf,
        ".docx": _extract_docx,
        ".doc":  _extract_doc,
        ".pptx": _extract_pptx,
        ".ppt":  _extract_doc,
        ".txt":  _extract_txt,
    }
    extractor = dispatch.get(ext, _extract_txt)
    return extractor(file_path)


# ─────────────────────────────────────────────────────────────────────────────
# Text preprocessing
# ─────────────────────────────────────────────────────────────────────────────

def preprocess_text(raw: str) -> str:
    """
    Clean extracted text for LLM consumption:
      1. Normalise Unicode to NFC (consistent encoding).
      2. Strip C0/C1 control characters except tab, LF, CR.
      3. Collapse horizontal whitespace (spaces/tabs) to single space.
      4. Collapse runs of 3+ blank lines to 2 blank lines.
      5. Strip leading/trailing whitespace from each line.
      6. Final strip of the whole string.
    """
    if not raw:
        return ""

    # 1 — NFC normalisation
    text = unicodedata.normalize("NFC", raw)

    # 2 — Strip control characters (keep \t, \n, \r)
    text = re.sub(r"[^\S\n\r\t ]+", " ", text)                      # replace odd whitespace with space
    text = "".join(
        ch for ch in text
        if unicodedata.category(ch) not in ("Cc", "Cf") or ch in "\n\r\t"
    )

    # 3 — Collapse horizontal whitespace
    text = re.sub(r"[ \t]+", " ", text)

    # 4 — Normalise line endings, then collapse excess blank lines
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"\n{3,}", "\n\n", text)

    # 5 — Strip per-line leading/trailing spaces
    lines = [line.strip() for line in text.split("\n")]
    text  = "\n".join(lines)

    return text.strip()


# ─────────────────────────────────────────────────────────────────────────────
# Keyword-relevance chunking helper
# ─────────────────────────────────────────────────────────────────────────────

def _score_paragraph(paragraph: str, keywords: set[str]) -> float:
    """Simple keyword hit-rate score for paragraph relevance."""
    if not keywords or not paragraph:
        return 0.0
    words = set(re.findall(r"\b\w+\b", paragraph.lower()))
    return len(words & keywords) / max(len(paragraph), 1) * 1000


def get_relevant_chunk(
    text: str,
    query: str = "",
    max_chars: int = MAX_CONTEXT_CHARS,
) -> str:
    """
    Return up to `max_chars` of `text`.

    When `query` is provided, paragraphs are ranked by keyword-overlap with
    the query and the top-scoring paragraphs are assembled first (up to the
    char budget), followed by any remaining space filled from the document
    top.  This gives query-relevant context priority without losing the
    document's overall structure entirely.

    When `query` is empty, a simple head-truncation is used.
    """
    if len(text) <= max_chars:
        return text

    if not query.strip():
        return text[:max_chars] + "\n\n[... document truncated for context window ...]"

    keywords = set(re.findall(r"\b\w{4,}\b", query.lower()))
    paragraphs = [p.strip() for p in re.split(r"\n{2,}", text) if p.strip()]

    scored = sorted(
        enumerate(paragraphs),
        key=lambda ip: _score_paragraph(ip[1], keywords),
        reverse=True,
    )

    selected_indices: set[int] = set()
    budget = max_chars - 100  # leave room for truncation note
    collected = 0

    for idx, para in scored:
        if collected + len(para) + 2 <= budget:
            selected_indices.add(idx)
            collected += len(para) + 2
        if collected >= budget:
            break

    # Preserve document order for selected paragraphs
    result_parts = [
        paragraphs[i] for i in sorted(selected_indices)
    ]
    result = "\n\n".join(result_parts)

    if len(result) < budget and paragraphs:
        result += "\n\n[... additional document sections omitted for context window ...]"

    return result


# ─────────────────────────────────────────────────────────────────────────────
# DocumentService
# ─────────────────────────────────────────────────────────────────────────────

class DocumentService:

    @staticmethod
    def ensure_upload_dir(user_id: int) -> Path:
        """Create the user upload directory if it does not exist."""
        user_dir = UPLOAD_DIR / str(user_id)
        user_dir.mkdir(parents=True, exist_ok=True)
        return user_dir

    @staticmethod
    def _find_duplicate(filename: str, file_size: int, user_id: int) -> Optional[str]:
        """
        Return the document_id if the user already has a document with
        the same filename and file size; else return None.
        """
        existing = DBService.list_documents(user_id)
        for doc in existing:
            if doc.get("filename") == filename and doc.get("file_size") == file_size:
                return doc["id"]
        return None

    @staticmethod
    async def save_upload(
        file_content: bytes,
        filename: str,
        user_id: int,
    ) -> dict:
        """
        Persist an uploaded file to disk, extract and preprocess its text,
        and store metadata in the database.

        Returns:
            {"document_id": str, "filename": str, "file_type": str,
             "file_size": int, "duplicate": bool}

        Raises ValueError on unsupported extension or oversized file.
        """
        # Validate extension
        ext = Path(filename).suffix.lower()
        if ext not in ALLOWED_EXTENSIONS:
            raise ValueError(
                f"Unsupported file type '{ext}'. "
                f"Allowed: {', '.join(sorted(ALLOWED_EXTENSIONS))}"
            )

        # Validate size
        if len(file_content) > MAX_FILE_SIZE_BYTES:
            raise ValueError(
                f"File too large ({len(file_content) // (1024 * 1024)} MB). "
                f"Maximum allowed size is {MAX_FILE_SIZE_BYTES // (1024 * 1024)} MB."
            )

        # Duplicate detection — return existing doc if same file re-uploaded
        duplicate_id = DocumentService._find_duplicate(filename, len(file_content), user_id)
        if duplicate_id:
            logger.info(
                "Duplicate upload detected for user %s — returning existing doc_id=%s",
                user_id, duplicate_id,
            )
            return {
                "document_id": duplicate_id,
                "filename": filename,
                "file_type": ext,
                "file_size": len(file_content),
                "duplicate": True,
            }

        # Generate IDs and paths
        doc_id    = str(uuid.uuid4())
        safe_name = f"{doc_id}_{Path(filename).name}"
        user_dir  = DocumentService.ensure_upload_dir(user_id)
        file_path = user_dir / safe_name

        # Write to disk (async)
        async with aiofiles.open(file_path, "wb") as f:
            await f.write(file_content)

        # Store metadata in DB
        DBService.create_document(
            doc_id=doc_id,
            user_id=user_id,
            filename=filename,
            file_path=str(file_path),
            file_type=ext,
            file_size=len(file_content),
        )

        # Extract + preprocess text in a thread (keeps event loop unblocked)
        try:
            raw_text = await asyncio.to_thread(extract_text, file_path, ext)
            if raw_text:
                clean_text = preprocess_text(raw_text)
                DBService.update_document_text(doc_id, clean_text)
                logger.info(
                    "Extracted & preprocessed %d chars from '%s' (doc_id=%s)",
                    len(clean_text), filename, doc_id,
                )
            else:
                logger.warning("No text extracted from '%s'", filename)
        except Exception as exc:
            logger.error("Text extraction error for '%s': %s", filename, exc)

        return {
            "document_id": doc_id,
            "filename": filename,
            "file_type": ext,
            "file_size": len(file_content),
            "duplicate": False,
        }

    @staticmethod
    def get_document_context(doc_id: str, user_id: int) -> Optional[str]:
        """
        Retrieve the extracted text for a document owned by the user.
        Returns a head-truncated string up to MAX_CONTEXT_CHARS.
        Returns None if the document is not found or has no extracted text.
        """
        doc = DBService.get_document(doc_id, user_id)
        if not doc:
            return None
        text = doc.get("extracted_text") or ""
        if not text:
            return None
        if len(text) > MAX_CONTEXT_CHARS:
            text = text[:MAX_CONTEXT_CHARS]
            text += "\n\n[... document truncated for context window ...]"
        return text

    @staticmethod
    def get_document_context_for_query(
        doc_id: str,
        user_id: int,
        query: str = "",
    ) -> Optional[str]:
        """
        Like get_document_context, but uses keyword-relevance ranking
        to prioritise paragraphs most relevant to `query`.
        """
        doc = DBService.get_document(doc_id, user_id)
        if not doc:
            return None
        text = doc.get("extracted_text") or ""
        if not text:
            return None
        return get_relevant_chunk(text, query=query, max_chars=MAX_CONTEXT_CHARS)

    @staticmethod
    def list_documents(user_id: int) -> list[dict]:
        """Return metadata for all documents uploaded by this user."""
        return DBService.list_documents(user_id)

    @staticmethod
    def get_document_meta(doc_id: str, user_id: int) -> Optional[dict]:
        """Return full document metadata row or None if not found/not owned."""
        return DBService.get_document(doc_id, user_id)

    @staticmethod
    def delete_document(doc_id: str, user_id: int) -> bool:
        """
        Delete a document from disk and from the database.
        Returns True if deleted, False if not found or not owned.
        """
        doc = DBService.get_document(doc_id, user_id)
        if not doc:
            return False

        file_path = Path(doc.get("file_path", ""))
        if file_path.exists():
            try:
                file_path.unlink()
            except Exception as exc:
                logger.error("Could not delete file %s: %s", file_path, exc)

        return DBService.delete_document(doc_id, user_id)

    @staticmethod
    def cleanup_old_documents(max_age_hours: int = 24) -> int:
        """
        Remove document files older than max_age_hours from disk.
        DB records are kept (soft-delete pattern).
        Returns the count of files deleted.
        """
        deleted = 0
        cutoff  = datetime.datetime.utcnow() - datetime.timedelta(hours=max_age_hours)
        if not UPLOAD_DIR.exists():
            return 0
        for user_dir in UPLOAD_DIR.iterdir():
            if not user_dir.is_dir():
                continue
            for fpath in user_dir.iterdir():
                if not fpath.is_file():
                    continue
                mtime = datetime.datetime.utcfromtimestamp(fpath.stat().st_mtime)
                if mtime < cutoff:
                    try:
                        fpath.unlink()
                        deleted += 1
                        logger.info("Cleanup: removed old file %s", fpath)
                    except Exception as exc:
                        logger.error("Cleanup error for %s: %s", fpath, exc)
        return deleted
